#!/usr/bin/env python3
"""
Kimi .pptd → PowerPoint .pptx конвертер (FullHD 1920x1080)
С поддержкой рендеринга иконок Phosphor в PNG через IconRenderer.
"""

import yaml
import zipfile
import tempfile
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import re
from tools.icon_renderer import IconRenderer
import matplotlib

matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# ═══════════════════════════════════════════════════════════
# НАСТРОЙКИ РАЗРЕШЕНИЯ
# ═══════════════════════════════════════════════════════════
SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080
SOURCE_WIDTH_PX = 1920
SOURCE_HEIGHT_PX = 1080
SCALE_X = 1.0
SCALE_Y = 1.0

SLIDE_WIDTH = Inches(SLIDE_WIDTH_PX / 96)
SLIDE_HEIGHT = Inches(SLIDE_HEIGHT_PX / 96)

ICONS_BASE_DIR = None


def px_to_inches_x(px):
    return Inches((px * SCALE_X) / 96)


def px_to_inches_y(py):
    return Inches((py * SCALE_Y) / 96)


# ═══════════════════════════════════════════════════════════
# ГЕНЕРАЦИЯ ГРАФИКОВ
# ═══════════════════════════════════════════════════════════

def generate_chart(chart_data, chart_type, colors, width=600, height=300, dpi=150):
    """Генерация графика с защитой от пустых/нулевых данных."""
    if not chart_data or not isinstance(chart_data, list):
        return None

    # Извлекаем метки и значения
    first_item = chart_data[0]
    if not isinstance(first_item, dict):
        return None

    keys = list(first_item.keys())
    if len(keys) < 2 and chart_type != 'pie':
        return None

    label_key = keys[0]
    value_key = keys[1] if len(keys) > 1 else None

    labels = [str(item.get(label_key, '')) for item in chart_data]

    values = []
    for item in chart_data:
        raw_val = item.get(value_key, 0) if value_key else item.get(keys[0], 0) if chart_type == 'pie' else 0
        try:
            val = float(raw_val)
        except (TypeError, ValueError):
            val = 0.0
        values.append(val)

    # Проверяем, есть ли осмысленные данные
    if all(v == 0 for v in values) and chart_type != 'pie':
        return None
    if sum(values) == 0 and chart_type == 'pie':
        return None

    fig_width = width / dpi
    fig_height = height / dpi
    fig, ax = plt.subplots(figsize=(fig_width, fig_height), dpi=dpi)
    fig.patch.set_facecolor('#F0EDEA')
    ax.set_facecolor('#F0EDEA')

    def hex_to_rgb_float(hex_color):
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i + 2], 16) / 255 for i in (0, 2, 4))

    rgb_colors = [hex_to_rgb_float(c) for c in colors]

    try:
        if chart_type == "bar":
            bars = ax.bar(labels, values, color=rgb_colors[0] if rgb_colors else '#C9892E',
                          width=0.6, edgecolor='none')
            max_val = max(values) if values else 1
            if max_val == 0:
                max_val = 1
            for bar, val in zip(bars, values):
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max_val * 0.02,
                        str(val), ha='center', va='bottom', fontsize=9, color='#1C2524')
            ax.set_ylim(0, max_val * 1.15)

        elif chart_type == "pie":
            while len(rgb_colors) < len(labels):
                rgb_colors.append((0.4, 0.4, 0.4))
            wedges, texts, autotexts = ax.pie(values, labels=labels, colors=rgb_colors[:len(labels)],
                                              autopct='%1.0f%%', startangle=90,
                                              textprops={'fontsize': 9, 'color': '#1C2524'})
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontsize(8)

        elif chart_type == "line":
            ax.plot(labels, values, color=rgb_colors[0] if rgb_colors else '#C9892E',
                    linewidth=2, marker='o', markersize=4, label=value_key)
            ax.legend(loc='upper right', fontsize=8, frameon=False)
            max_val = max(values) if values else 1
            if max_val == 0:
                max_val = 1
            ax.set_ylim(0, max_val * 1.2)

        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.spines['bottom'].set_color('#D5D0CA')
        ax.tick_params(colors='#6B7D7B', labelsize=8)
        ax.grid(axis='y', alpha=0.3, color='#D5D0CA')

        tmp_path = Path(tempfile.gettempdir()) / f"chart_{chart_type}_{id(chart_data)}.png"
        plt.tight_layout(pad=0.3)
        plt.savefig(tmp_path, dpi=dpi, facecolor='#F0EDEA', edgecolor='none',
                    bbox_inches='tight', pad_inches=0.05)
        plt.close()
        return tmp_path
    except Exception as e:
        plt.close()
        print(f"⚠ Ошибка при построении графика: {e}")
        return None


# ═══════════════════════════════════════════════════════════
# КОНВЕРТЕР
# ═══════════════════════════════════════════════════════════

class KimiPptdConverter:
    def __init__(self, icons_base_dir=None):
        self.theme = {}
        self.colors = {}
        self.text_styles = {}
        self.table_styles = {}
        self.images_dir = None
        self.pages_dir = None
        self.icons_dir = icons_base_dir
        self._icon_cache = {}
        # Новый рендерер иконок
        self.icon_renderer = None
        if icons_base_dir and Path(icons_base_dir).exists():
            from tools.icon_renderer import IconRenderer
            self.icon_renderer = IconRenderer(Path(icons_base_dir))

    def load_project(self, project_dir):
        project_dir = Path(project_dir)
        pptd_files = list(project_dir.glob("*.pptd"))
        if not pptd_files:
            raise ValueError("Файл .pptd не найден в архиве")

        with open(pptd_files[0], 'r', encoding='utf-8') as f:
            pptd = yaml.safe_load(f)

        self.theme = pptd.get("theme", {})
        self.colors = self.theme.get("colors", {})
        self.text_styles = self.theme.get("textStyles", {})

        # === НОВОЕ: fontPair ===
        font_pair = self.theme.get("fontPair", "Arial + Arial")
        self.title_font, self.body_font = self._parse_font_pair(font_pair)

        # Применяем fontPair к text_styles если там не указаны шрифты
        for style_name, style in self.text_styles.items():
            if "fontFamily" not in style:
                if style_name in ("title", "subtitle"):
                    style["fontFamily"] = self.title_font
                else:
                    style["fontFamily"] = self.body_font
        self.table_styles = self.theme.get("tableStyles", {})

        self.pages_dir = project_dir / "pages"
        self.images_dir = project_dir / "images"

        return pptd.get("pages", [])

    def _parse_font_pair(self, font_pair: str) -> tuple:
        """'Montserrat + Inter' → ('Montserrat', 'Inter')"""
        parts = font_pair.replace(" + ", "+").split("+")
        if len(parts) == 2:
            return parts[0].strip(), parts[1].strip()
        return "Arial", "Arial"

    def resolve_color(self, val):
        if isinstance(val, str) and val.startswith("$"):
            return self.colors.get(val[1:], "#000000")
        return val

    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))

    def strip_html_tags(self, text_str):
        if not text_str:
            return ""
        try:
            root = etree.fromstring(f"<root>{text_str}</root>")
            return etree.tostring(root, method='text', encoding='unicode')
        except:
            clean = re.sub(r'<[^>]+>', '', text_str)
            return clean

    def parse_html_text(self, text_str):
        text_str = text_str.strip()
        if not text_str:
            return [("", False, False, None)]

        try:
            root = etree.fromstring(f"<root>{text_str}</root>")
        except:
            clean = self.strip_html_tags(text_str)
            return [(clean, False, False, None)]

        result = []

        def walk(node, bold=False, italic=False, color=None):
            if node.tag in ("strong", "b"):
                bold = True
            if node.tag in ("em", "i"):
                italic = True

            if 'style' in node.attrib:
                style = node.attrib['style']
                color_match = re.search(r'color:\s*([^;\s"]+)', style)
                if color_match:
                    color = color_match.group(1).strip()

            if node.text and node.text.strip():
                result.append((node.text.strip(), bold, italic, color))

            for child in node:
                walk(child, bold, italic, color)
                if child.tail and child.tail.strip():
                    result.append((child.tail.strip(), bold, italic, color))

        walk(root)
        return result if result else [("", False, False, None)]

    def convert(self, project_dir, output_path, progress_callback=None):
        page_refs = self.load_project(project_dir)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        slide_layout = prs.slide_layouts[6]

        total = len(page_refs)

        for i, page_ref in enumerate(page_refs):
            if progress_callback:
                progress_callback(int((i / total) * 100))

            page_path = self.pages_dir / Path(page_ref).name
            if not page_path.exists():
                print(f"⚠ Пропущен: {page_ref}")
                continue

            with open(page_path, 'r', encoding='utf-8') as f:
                page = yaml.safe_load(f)

            slide = prs.slides.add_slide(slide_layout)

            # Фон
            bg = page.get("background", {})
            if bg.get("type") == "solid":
                bg_color = self.resolve_color(bg.get("color", "$white"))
                r, g, b = self.hex_to_rgb(bg_color)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)
            elif bg.get("type") == "image":
                img_name = bg.get("image", "")
                img_path = self.images_dir / img_name if img_name else None
                if img_path and img_path.exists():
                    slide.shapes.add_picture(
                        str(img_path), Inches(0), Inches(0),
                        SLIDE_WIDTH, SLIDE_HEIGHT
                    )

            # Элементы
            for el in page.get("elements", []):
                self._process_element(slide, el)

            if progress_callback:
                progress_callback(int(((i + 1) / total) * 100))

        prs.save(output_path)
        return output_path

    def _process_element(self, slide, el):
        el_type = el.get("elementType")
        bounds = el.get("bounds", [0, 0, 100, 100])
        x, y, w, h = bounds

        if el_type == "shape":
            self._add_shape(slide, el, x, y, w, h)
        elif el_type == "text":
            self._add_text(slide, el, x, y, w, h)
        elif el_type == "image":
            self._add_image(slide, el, x, y, w, h)
        elif el_type == "table":
            self._add_table(slide, el, x, y, w, h)
        elif el_type == "icon":
            self._add_icon(slide, el, x, y, w, h)
        elif el_type == "chart":
            self._add_chart(slide, el, x, y, w, h)

    def _add_shape(self, slide, el, x, y, w, h):
        shape_name = el.get("shapeName", "rect")
        shape_map = {
            "rect": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "oval": MSO_SHAPE.OVAL,
            "roundedRect": MSO_SHAPE.ROUNDED_RECTANGLE,
        }
        shape_type = shape_map.get(shape_name, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(
            shape_type,
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        )

        fill_cfg = el.get("fill", {})
        if fill_cfg.get("type") == "solid":
            fill_color = self.resolve_color(fill_cfg.get("color", "$white"))
            r, g, b = self.hex_to_rgb(fill_color)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        elif fill_cfg.get("type") == "gradient":
            stops = fill_cfg.get("stops", [])
            if stops:
                fill_color = self.resolve_color(stops[0].get("color", "$white"))
                r, g, b = self.hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            shape.fill.background()

        border = el.get("border", {})
        if border:
            border_color = self.resolve_color(border.get("color", "$white"))
            r, g, b = self.hex_to_rgb(border_color)
            shape.line.color.rgb = RGBColor(r, g, b)
            shape.line.width = Pt(border.get("width", 1))
        else:
            shape.line.fill.background()

    def _add_text(self, slide, el, x, y, w, h):
        content = el.get("content", {})
        style_ref = content.get("style", "")

        style = {}
        if style_ref.startswith("$"):
            style = self.text_styles.get(style_ref[1:], {})

        # === НОВОЕ: spacing scale для line_height ===
        spacing_scale = self.theme.get("spacingScale", [8, 16, 24, 32, 48, 64])
        base_line_height = style.get("lineHeight", 1.2)

        # Округляем line_height до ближайшего из spacing_scale / fontSize
        font_size = content.get("fontSize", style.get("fontSize", 18))
        ideal_spacing = font_size * base_line_height
        snapped_spacing = min(spacing_scale, key=lambda s: abs(s - ideal_spacing))
        line_height = snapped_spacing / font_size if font_size > 0 else 1.2
        align_cfg = content.get("align", ["left", "top"])
        h_align = align_cfg[0] if len(align_cfg) > 0 else "left"

        txBox = slide.shapes.add_textbox(
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        )
        tf = txBox.text_frame
        tf.word_wrap = content.get("wrap", True)
        tf.clear()

        text_raw = content.get("text", "")
        paragraphs = [p.strip() for p in text_raw.strip().split("\n") if p.strip()]

        # Автоподбор размера шрифта
        text_raw = content.get("text", "")
        style_info = self.text_styles.get(style_ref.lstrip("$"), {})
        base_size = style_info.get("fontSize", 18)

        # Оцениваем, влезет ли текст
        estimated_lines = len(text_raw) / (w / (base_size * 0.6))  # примерная ширина символа
        available_lines = h / (base_size * 1.5)

        if estimated_lines > available_lines * 1.2:
            # Уменьшаем шрифт пропорционально
            scale = (available_lines * 1.2) / estimated_lines
            font_size = max(14, int(base_size * scale))
        else:
            font_size = base_size

        for idx, para_text in enumerate(paragraphs):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()

            if h_align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif h_align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            parsed = self.parse_html_text(para_text)
            for text_part, is_bold, is_italic, inline_color in parsed:
                if not text_part:
                    continue
                run = p.add_run()
                run.text = text_part
                run.font.size = Pt(font_size)
                run.font.name = font_family
                run.font.bold = is_bold
                run.font.italic = is_italic

                final_color = inline_color if inline_color else base_color
                r, g, b = self.hex_to_rgb(self.resolve_color(final_color))
                run.font.color.rgb = RGBColor(r, g, b)

            p.space_after = Pt(0)
            p.space_before = Pt(0)
            if line_height > 1.0:
                p.line_spacing = line_height

    def _add_image(self, slide, el, x, y, w, h):
        img_name = el.get("image", "")
        src = el.get("src", "")

        img_path = None

        if img_name:
            local_path = self.images_dir / img_name
            if local_path.exists():
                img_path = local_path
        elif src:
            url_name = src.split("/")[-1].split("?")[0]
            if url_name:
                local_path = self.images_dir / url_name
                if local_path.exists():
                    img_path = local_path

        if img_path and img_path.exists():
            try:
                slide.shapes.add_picture(
                    str(img_path),
                    px_to_inches_x(x), px_to_inches_y(y),
                    px_to_inches_x(w), px_to_inches_y(h)
                )
            except Exception as e:
                print(f"⚠ Ошибка добавления изображения {img_path.name}: {e}")

    def _add_table(self, slide, el, x, y, w, h):
        rows = el.get("rows", [])
        if not rows:
            return

        num_rows = len(rows)
        num_cols = max(len(row.get("cells", [])) for row in rows)

        table = slide.shapes.add_table(
            num_rows, num_cols,
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        ).table

        style = self.table_styles.get("default", {})
        header_fill = self.resolve_color(style.get("headerFill", "$primary"))
        header_color = self.resolve_color(style.get("headerColor", "$textLight"))
        body_fill_colors = style.get("bodyFill", ["$bgLight", "$cardBg"])
        body_color = self.resolve_color(style.get("bodyColor", "$textDark"))

        for i, row in enumerate(rows):
            cells = row.get("cells", [])
            for j, cell_text in enumerate(cells):
                if j >= num_cols:
                    break
                cell = table.cell(i, j)
                cell.text = str(cell_text)

                if i == 0 and row.get("isHeader", False):
                    r, g, b = self.hex_to_rgb(header_fill)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(r, g, b)
                    text_color = header_color
                else:
                    fill_idx = (i - (1 if rows[0].get("isHeader", False) else 0)) % len(body_fill_colors)
                    fill_color = self.resolve_color(body_fill_colors[fill_idx])
                    r, g, b = self.hex_to_rgb(fill_color)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(r, g, b)
                    text_color = body_color

                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        r, g, b = self.hex_to_rgb(text_color)
                        run.font.color.rgb = RGBColor(r, g, b)
                        run.font.size = Pt(style.get("fontSize", 16))
                        run.font.name = style.get("fontFamily", "Arial")


    def _add_icon(self, slide, el, x, y, w, h):
        """Вставляет иконку как PNG с прозрачностью, используя IconRenderer."""
        icon_name = el.get("iconName", "")
        if isinstance(icon_name, dict):
            icon_name = icon_name.get("name", "circle")
        elif not isinstance(icon_name, str):
            icon_name = str(icon_name)
        fill_cfg = el.get("fill", {})
        icon_color = self.resolve_color(fill_cfg.get("color", "$accent"))
        style = el.get("style", "regular")
        decorative = el.get("decorativeShape")
        shape_color = self.resolve_color(el.get("shapeColor", "$accent")) if decorative else None

        if self.icon_renderer:
            try:
                png_path = self.icon_renderer.render(
                    icon_name=icon_name,
                    size=max(w, h),
                    fill_color=icon_color,
                    style=style,
                    decorative_shape=decorative,
                    shape_color=shape_color,
                )
                slide.shapes.add_picture(
                    str(png_path),
                    px_to_inches_x(x), px_to_inches_y(y),
                    px_to_inches_x(w), px_to_inches_y(h)
                )
                return
            except Exception as e:
                print(f"⚠ Ошибка рендеринга иконки {icon_name}: {e}")

        # Fallback — круг с буквами
        self._draw_fallback_icon(slide, x, y, w, h, self.hex_to_rgb(icon_color), icon_name)

    def _draw_fallback_icon(self, slide, x, y, w, h, color, icon_name):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon_name.replace("fas:", "").replace("ph:", "")[:2].upper()
        font_size = max(8, min(w, h) * 0.4)
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True

    def _add_chart(self, slide, el, x, y, w, h):
        chart_type = el.get("type", "bar")
        data = el.get("data", [])
        colors_cfg = el.get("colors", ["$accent"])

        if not data:
            return

        colors = [self.resolve_color(c) for c in colors_cfg]

        try:
            chart_path = generate_chart(
                data, chart_type, colors,
                width=int(w * SCALE_X),
                height=int(h * SCALE_Y),
                dpi=150
            )

            if chart_path and chart_path.exists():
                slide.shapes.add_picture(
                    str(chart_path),
                    px_to_inches_x(x), px_to_inches_y(y),
                    px_to_inches_x(w), px_to_inches_y(h)
                )
        except Exception as e:
            print(f"⚠ Ошибка генерации графика: {e}")
            self._add_chart_fallback(slide, el, x, y, w, h)

    def _add_chart_fallback(self, slide, el, x, y, w, h):
        chart_type = el.get("type", "chart")
        data = el.get("data", [])

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 237, 234)
        shape.line.color.rgb = RGBColor(213, 208, 202)
        shape.line.width = Pt(1)

        txBox = slide.shapes.add_textbox(
            px_to_inches_x(x), px_to_inches_y(y),
            px_to_inches_x(w), px_to_inches_y(h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        chart_labels = {
            "bar": "📊 Гистограмма",
            "pie": "🥧 Круговая диаграмма",
            "line": "📈 Линейный график"
        }
        label = chart_labels.get(chart_type, f"📊 {chart_type}")

        run = p.add_run()
        run.text = f"{label}\n({len(data)} точек данных)"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(107, 125, 123)


def convert_project_to_pptx(zip_path, output_path=None, progress_callback=None, icons_base_dir=None):
    zip_path = Path(zip_path)
    if output_path is None:
        output_path = zip_path.with_suffix(".pptx")

    with tempfile.TemporaryDirectory() as tmpdir:
        with zipfile.ZipFile(zip_path, 'r') as z:
            z.extractall(tmpdir)

        project_dir = Path(tmpdir)
        subdirs = [d for d in project_dir.iterdir() if d.is_dir()]
        if subdirs and not list(project_dir.glob("*.pptd")):
            project_dir = subdirs[0]

        converter = KimiPptdConverter(icons_base_dir=icons_base_dir)
        converter.convert(project_dir, output_path, progress_callback)

    return output_path